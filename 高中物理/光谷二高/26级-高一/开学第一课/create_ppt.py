from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE

OUT = '高中物理/光谷二高/26级-高一/开学第一课/开学第一课_简洁版_v2.pptx'
FONT = 'Microsoft YaHei'
NAVY = RGBColor(9, 20, 38)
INK = RGBColor(24, 39, 62)
MUTED = RGBColor(103, 122, 148)
CYAN = RGBColor(38, 197, 220)
YELLOW = RGBColor(255, 194, 71)
WHITE = RGBColor(248, 251, 253)
PALE = RGBColor(235, 245, 248)
LAV = RGBColor(239, 237, 250)
MINT = RGBColor(229, 247, 240)
CORAL = RGBColor(255, 232, 224)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box

def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.12
    return shape

def circle(slide, x, y, d, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape

def line(slide, x1, y1, x2, y2, color=CYAN, width=2.2, dash=None):
    sh = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    sh.line.color.rgb = color; sh.line.width = Pt(width)
    if dash: sh.line.dash_style = dash
    return sh

def arrow(slide, x1, y1, x2, y2, color=CYAN, width=2.5):
    sh = line(slide,x1,y1,x2,y2,color,width)
    sh.line.end_arrowhead = True
    return sh

def header(slide, kicker, title, num):
    add_text(slide, kicker.upper(), .65, .35, 3.0, .28, 10, CYAN, True)
    add_text(slide, title, .65, .68, 11.4, .55, 27, NAVY, True)
    add_text(slide, f'{num:02d}  /  开学第一课', 11.25, .42, 1.4, .25, 9, MUTED, False, PP_ALIGN.RIGHT)
    line(slide, .65, 1.35, 12.65, 1.35, RGBColor(210,224,232), 1)

def footer(slide, text='高中物理 · 光谷二高'):
    add_text(slide, text, .68, 7.12, 4.0, .2, 9, MUTED)

def add_bullet(slide, x, y, label, desc, fill):
    circle(slide,x,y+.08,.25,fill)
    add_text(slide,label,x+.42,y,2.1,.3,17,INK,True)
    add_text(slide,desc,x+.42,y+.34,3.2,.45,11,MUTED)

# 1 cover
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,NAVY)
for i,(x,y,d,c) in enumerate([(9.3,1.0,2.5,CYAN),(10.7,2.2,1.2,YELLOW),(8.7,4.2,1.0,WHITE)]):
    circle(s,x,y,d,c)
for a,b,c,d in [(7.0,5.8,12.0,1.6),(7.2,6.1,12.3,3.0),(8.0,5.6,11.7,5.7)]: line(s,a,b,c,d,RGBColor(65,104,129),1.6)
add_text(s,'开学第一课',.8,1.02,6.8,.8,42,WHITE,True)
add_text(s,'从“我不会物理”\n到“我能解释现象”',.82,2.05,6.8,1.35,27,CYAN,True)
add_text(s,'高一新生 · 40分钟',.85,4.15,4.0,.35,15,RGBColor(190,213,224))
add_text(s,'看清现象  ·  说出理由  ·  写完整过程',.85,5.05,6.5,.35,14,WHITE)
add_text(s,'刘老师  |  2026',.85,6.55,4.0,.3,11,RGBColor(157,183,198)); add_text(s,'目录 01 · 开场与约定',9.2,6.55,3.0,.25,10,RGBColor(157,183,198),False,PP_ALIGN.RIGHT)

# 2 directory
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'课程目录','今天，我们沿着五个部分认识高中物理',2)
items=[('01','开场与约定','认识彼此，建立安全的课堂氛围'),('02','物理从哪里开始','从生活现象进入物理问题'),('03','初中基础 → 高中地图','先回顾，再看高中主线'),('04','学习方法','费曼法 + 三分钟小老师'),('05','高一上任务与行动','数学工具、习惯与出口条')]
for i,(n,t,d) in enumerate(items):
    y=1.72+i*.88; circle(s,.95,y,.48,[CYAN,YELLOW,RGBColor(135,202,171),RGBColor(188,166,230),RGBColor(255,161,123)][i]); add_text(s,n,.95,y+.1,.48,.22,10,NAVY,True,PP_ALIGN.CENTER); add_text(s,t,1.7,y-.02,3.7,.3,17,INK,True); add_text(s,d,5.45,y-.02,6.5,.3,13,MUTED)
footer(s,'目录总览 · 5 个部分 / 14 页')

# 3 agreements
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'先认识彼此','三条课堂约定，让“不会”成为起点',3)
add_bullet(s,.9,2.0,'允许不会','不会是起点，不是结论。',CYAN)
add_bullet(s,.9,3.35,'先想再问','先说现象、办法，再说卡点。',YELLOW)
add_bullet(s,.9,4.7,'用证据说话','答案要有依据、过程和单位。',RGBColor(135,202,171))
rect(s,7.0,1.95,5.2,3.8,PALE,True)
add_text(s,'今天请写下',7.45,2.3,4.2,.4,18,NAVY,True)
for i,t in enumerate(['姓名','初中最熟悉的内容','对高中物理最担心的事']):
    circle(s,7.5,3.05+i*.75,.3,[CYAN,YELLOW,RGBColor(135,202,171)][i])
    add_text(s,t,7.95,3.0+i*.75,3.6,.35,15,INK)
footer(s,'目录 01 · 开场与约定')

# 4 phenomena
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'破冰','物理，藏在每一个“为什么”里',4)
cards=[('启动 / 刹车','人为什么会前仰后仰？',CYAN,'↗'),('滚动距离','地板和草地，哪里停得更快？',YELLOW,'●'),('电梯启动','为什么刚启动会有“失重感”？',RGBColor(135,202,171),'↕')]
for i,(t,d,c,ic) in enumerate(cards):
    x=.85+i*4.12; rect(s,x,2.05,3.55,3.8,PALE if i==0 else (RGBColor(255,247,228) if i==1 else MINT),True)
    circle(s,x+.28,2.35,.75,c); add_text(s,ic,x+.28,2.37,.75,.55,26,NAVY,True,PP_ALIGN.CENTER)
    add_text(s,t,x+.28,3.35,2.9,.35,19,INK,True); add_text(s,d,x+.28,3.95,2.9,.8,15,MUTED)
    add_text(s,'先提出问题，不急着回答',x+.28,5.18,2.9,.3,11,c,True)
footer(s,'目录 02 · 物理从哪里开始')

# 5 route
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,NAVY); add_text(s,'学习路线',.7,.45,3,.3,10,CYAN,True); add_text(s,'高中物理学习 = 一条可重复的路线',.7,.78,10,.6,29,WHITE,True); add_text(s,'每一步都能留下证据',.7,1.38,5,.3,13,RGBColor(183,211,222)); add_text(s,'05  /  开学第一课',11.25,.42,1.4,.25,9,RGBColor(150,182,197),False,PP_ALIGN.RIGHT)
steps=[('观察','看到什么'),('问题','想知道什么'),('模型','先忽略什么'),('规律','关系是什么'),('证据','如何检验'),('表达','能否讲清')]
for i,(t,d) in enumerate(steps):
    x=.85+i*2.05; circle(s,x,3.0,.9,[CYAN,YELLOW,RGBColor(160,216,184),RGBColor(190,170,235),RGBColor(255,159,123),WHITE][i])
    add_text(s,str(i+1),x,3.16,.9,.4,19,NAVY,True,PP_ALIGN.CENTER); add_text(s,t,x-.2,4.15,1.3,.3,16,WHITE,True,PP_ALIGN.CENTER); add_text(s,d,x-.45,4.55,1.8,.3,11,RGBColor(177,204,217),False,PP_ALIGN.CENTER)
    if i<5: arrow(s,x+.95,3.45,x+1.92,3.45,RGBColor(93,132,151),2)
add_text(s,'遇到难题，就把它拆回上一步。',.8,6.2,6,.4,18,YELLOW,True); footer(s,'目录 02 · 物理从哪里开始')

# 6 junior map
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'知识回顾','先回顾初中物理：五个熟悉的领域',6)
jmods=[('机械运动','位置 · 路程 · 时间 · 速度',CYAN),('力与运动','重力 · 弹力 · 摩擦力',YELLOW),('声光热','传播 · 反射折射 · 温度热量',RGBColor(135,202,171)),('电与磁','电路 · 电流电压电阻',RGBColor(188,166,230)),('能量','功 · 功率 · 机械能转化',RGBColor(255,161,123))]
for i,(t,d,c) in enumerate(jmods):
    col=i%3; row=i//3; x=.95+col*4.1; y=1.9+row*1.85; rect(s,x,y,3.5,1.25,RGBColor(247,250,252),True,line=RGBColor(220,231,237)); rect(s,x,y,.12,1.25,c); add_text(s,f'{i+1:02d}',x+.3,y+.2,.45,.25,11,c,True); add_text(s,t,x+.3,y+.5,2.9,.3,17,INK,True); add_text(s,d,x+.3,y+.87,2.9,.25,11,MUTED)
add_text(s,'这些内容回答“现象是什么、规律怎样用”。',.95,6.25,11,.35,16,NAVY,True,PP_ALIGN.CENTER); footer(s,'目录 03 · 初中基础 → 高中地图')

# 7 senior map
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'知识地图','再看高中物理：沿着六条主线展开',7)
mods=[('运动的描述','位移 · 速度 · 加速度',CYAN),('相互作用','力 · 牛顿定律',YELLOW),('曲线与引力','平抛 · 圆周 · 天体',RGBColor(135,202,171)),('能量与动量','功 · 守恒 · 碰撞',RGBColor(188,166,230)),('电场与电路','电荷 · 电势 · 电流',RGBColor(255,161,123)),('磁场与近代','电磁感应 · 波 · 原子',RGBColor(131,173,223))]
for i,(t,d,c) in enumerate(mods):
    row=i//3; col=i%3; x=.85+col*4.1; y=1.9+row*2.05
    rect(s,x,y,3.55,1.45,RGBColor(247,250,252),True,line=RGBColor(220,231,237)); rect(s,x,y,.12,1.45,c)
    add_text(s,f'{i+1:02d}',x+.3,y+.22,.45,.25,11,c,True); add_text(s,t,x+.3,y+.53,2.9,.3,17,INK,True); add_text(s,d,x+.3,y+.92,2.9,.28,11,MUTED)
footer(s,'目录 03 · 初中基础 → 高中地图')

# 8 compare
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'台阶变高','初中与高中：不是换语言，而是多了模型',8)
add_text(s,'初中',1.2,1.8,3,.4,20,MUTED,True,PP_ALIGN.CENTER); add_text(s,'高中',8.2,1.8,3,.4,20,CYAN,True,PP_ALIGN.CENTER)
rows=[('现象直观','先建模，再分析'),('代入公式','函数、方程、向量'),('条件直接给出','从题干提取条件'),('看结果','查过程、单位、意义')]
for i,(a,b) in enumerate(rows):
    y=2.45+i*.83; rect(s,1.0,y,3.3,.55,PALE,True); rect(s,8.0,y,3.3,.55,RGBColor(232,248,250),True)
    add_text(s,a,1.15,y+.05,3,.3,15,INK,True,PP_ALIGN.CENTER); add_text(s,b,8.15,y+.05,3,.3,15,INK,True,PP_ALIGN.CENTER); arrow(s,4.65,y+.28,7.7,y+.28,MUTED,1.5)
add_text(s,'数学暂时薄弱可以补；关键是愿意画图、写步骤、问清每个量。',1.2,6.25,10.8,.4,16,NAVY,True,PP_ALIGN.CENTER); footer(s,'目录 03 · 初中基础 → 高中地图')

# 9 feynman
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'核心方法','费曼学习法：把“会做”变成“会讲”',9)
cx,cy=6.65,3.75; circle(s,cx-.72,cy-.72,1.44,NAVY); add_text(s,'真正\n学会',cx-.65,cy-.42,1.3,.7,18,WHITE,True,PP_ALIGN.CENTER)
fsteps=[('1','选小概念','只讲一个点',CYAN,2.1,2.0),('2','用自己的话','像讲给小学生',YELLOW,8.9,2.0),('3','找出卡点','回看资料',RGBColor(135,202,171),2.1,5.0),('4','简化举例','再讲一遍',RGBColor(188,166,230),8.9,5.0)]
for n,t,d,c,x,y in fsteps:
    circle(s,x,y,.62,c); add_text(s,n,x,y+.12,.62,.3,16,NAVY,True,PP_ALIGN.CENTER); add_text(s,t,x+.82,y-.02,2.9,.3,16,INK,True); add_text(s,d,x+.82,y+.32,2.9,.3,12,MUTED)
arrow(s,3.05,2.3,5.85,3.1,MUTED,1.4); arrow(s,7.45,3.1,9.0,2.3,MUTED,1.4); arrow(s,9.0,5.45,7.45,4.4,MUTED,1.4); arrow(s,5.85,4.4,3.05,5.45,MUTED,1.4); footer(s,'目录 04 · 学习方法｜卡住不是失败，卡点就是下一步')

# 10 concept card
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'示范','概念卡：用“速度”练习说清楚',10)
rect(s,.85,1.85,6.0,4.65,RGBColor(250,252,253),True,line=RGBColor(210,225,232)); add_text(s,'概念讲解卡',1.2,2.15,2.5,.35,20,NAVY,True)
fields=[('我解释的概念','速度'),('用自己的话说','单位时间内位移的变化快慢'),('生活例子','2 s 向东 10 m，再 2 s 向西 10 m'),('容易混淆','速率：路程/时间；速度：位移/时间')]
for i,(k,v) in enumerate(fields):
    y=2.8+i*.78; add_text(s,k,1.2,y,1.6,.25,11,MUTED,True); add_text(s,v,2.85,y,3.45,.38,14,INK, i==0)
    line(s,2.85,y+.45,6.35,y+.45,RGBColor(220,231,237),1)
rect(s,7.5,2.0,4.6,3.8,PALE,True); add_text(s,'固定表达模板',7.9,2.35,3.6,.35,18,NAVY,True)
add_text(s,'我研究的是____\n它表示____\n例如____\n容易和____混淆，区别是____',7.9,3.05,3.6,1.8,18,INK)
add_text(s,'30 秒讲清：是什么 · 例子 · 易错点',7.9,5.2,3.7,.3,12,CYAN,True); footer(s,'目录 04 · 学习方法')

# 11 group activity
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'课堂练习','三分钟小老师：准备—讲解—追问',11)
for i,(t,d,c) in enumerate([('1 分钟','分工列提纲',CYAN),('1 分钟','一人讲解',YELLOW),('1 分钟','接受追问',RGBColor(135,202,171))]):
    x=1.0+i*4.05; circle(s,x,2.05,1.0,c); add_text(s,str(i+1),x,2.28,1,.35,22,NAVY,True,PP_ALIGN.CENTER); add_text(s,t,x-.25,3.35,1.5,.3,15,INK,True,PP_ALIGN.CENTER); add_text(s,d,x-.8,3.78,2.6,.35,13,MUTED,False,PP_ALIGN.CENTER)
    if i<2: arrow(s,x+1.1,2.55,x+3.45,2.55,MUTED,1.5)
rect(s,1.05,5.0,11.1,1.0,RGBColor(248,250,252),True,line=RGBColor(220,231,237)); add_text(s,'练习示例：解释“惯性”——公交车刹车时身体前倾；易错点：惯性不是力。',1.35,5.12,10.5,.55,15,NAVY,True,PP_ALIGN.CENTER); footer(s,'目录 04 · 学习方法｜主题：速度 / 惯性 / 摩擦力 / 电流')

# 12 tasks
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'本学期','高一上：先完成两件重要的事',12)
rect(s,.85,1.9,5.4,3.8,PALE,True); add_text(s,'任务 01',1.25,2.25,1.3,.25,11,CYAN,True); add_text(s,'思维方式转变',1.25,2.72,3.8,.4,22,NAVY,True); add_text(s,'从“等答案”\n到“先提出模型和理由”',1.25,3.5,4.1,1.0,19,INK)
rect(s,7.05,1.9,5.4,3.8,RGBColor(255,248,232),True); add_text(s,'任务 02',7.45,2.25,1.3,.25,11,YELLOW,True); add_text(s,'补齐够用的数学工具',7.45,2.72,4.5,.4,22,NAVY,True)
tools=['正负数、分数与公式变形','单位换算与数量级估算','一次方程：从文字列式','坐标、图像与斜率直觉','标量/矢量与同一直线加减','直角三角形求水平/竖直分量']
for i,t in enumerate(tools):
    circle(s,7.5+(i%2)*2.35,3.42+(i//2)*.62,.24,YELLOW); add_text(s,t,7.85+(i%2)*2.35,3.35+(i//2)*.62,1.95,.42,11,INK,True)
add_text(s,'每周 1–2 个小专题；学一点，用一次，再讲给别人听。',1.0,6.25,11.3,.35,16,NAVY,True,PP_ALIGN.CENTER); footer(s,'目录 05 · 高一上任务与行动')

# 13 habits and exit
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,WHITE); header(s,'带走行动','从本周开始，固定一个小动作',13)
hab=[('课前','10 分钟：读标题，圈 1 个不懂词'),('课中','先听思路，再记关键式子'),('课后','少量基础题 + 一次错题复盘')]
for i,(t,d) in enumerate(hab):
    x=.9+i*4.1; rect(s,x,1.95,3.45,1.55,[PALE,RGBColor(255,248,232),MINT][i],True); add_text(s,t,x+.25,2.2,1.0,.3,17,NAVY,True); add_text(s,d,x+.25,2.75,2.9,.45,12,INK)
rect(s,1.0,4.4,11.1,1.35,NAVY,True); add_text(s,'出口条',1.35,4.68,1.1,.3,16,CYAN,True); add_text(s,'今天我知道了：____    我仍然困惑：____    本周我坚持：____',2.55,4.66,8.9,.35,16,WHITE)
add_text(s,'请选择一项，写在你的本子上。',1.0,6.25,11.2,.35,16,NAVY,True,PP_ALIGN.CENTER); footer(s,'目录 05 · 高一上任务与行动')

# 14 closing
s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,NAVY)
circle(s,10.7,1.2,1.7,CYAN); circle(s,11.7,2.6,.65,YELLOW); line(s,7.2,5.7,12.1,2.3,RGBColor(69,111,137),1.5); line(s,7.2,5.7,11.1,5.9,RGBColor(69,111,137),1.5)
add_text(s,'从今天开始',.85,1.15,5.5,.55,35,WHITE,True); add_text(s,'看清现象\n说出理由\n写完整过程',.9,2.25,5.7,2.15,29,CYAN,True)
add_text(s,'每周多懂一个概念，多改正一个错误，\n你就在进步。',.92,5.25,7.0,.8,19,WHITE)
add_text(s,'欢迎来到物理的世界。',.92,6.55,5.0,.35,15,YELLOW,True); add_text(s,'14  /  开学第一课',11.25,.42,1.4,.25,9,RGBColor(150,182,197),False,PP_ALIGN.RIGHT); footer(s,'目录 05 · 高一上任务与行动')

prs.save(OUT)
print(OUT)
