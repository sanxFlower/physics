from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os, math

OUT = os.path.join(os.path.dirname(__file__), '1.2时间位移_开学第一课.pptx')
ASSET = os.path.join(os.path.dirname(__file__), '教学设计插图')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(17, 34, 64)
BLUE = RGBColor(29, 111, 191)
CYAN = RGBColor(45, 183, 199)
ORANGE = RGBColor(242, 146, 61)
YELLOW = RGBColor(247, 203, 74)
GREEN = RGBColor(55, 166, 116)
RED = RGBColor(214, 76, 76)
INK = RGBColor(33, 44, 58)
MUTED = RGBColor(100, 116, 133)
PALE = RGBColor(240, 246, 251)
WHITE = RGBColor(255, 255, 255)
GRID = RGBColor(218, 228, 238)

def set_bg(slide, color=WHITE):
    shape = slide.background
    fill = shape.fill
    fill.solid(); fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, line=None, radius=False, transparency=0):
    typ = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.fill.transparency = transparency
    s.line.color.rgb = line if line else fill
    if radius:
        s.adjustments[0] = 0.08
    return s

def line(slide, x1, y1, x2, y2, color=INK, width=2, dash=None, begin=None, end=None):
    l = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    l.line.color.rgb = color; l.line.width = Pt(width)
    if dash: l.line.dash_style = dash
    if begin: l.line.begin_arrowhead = begin
    if end: l.line.end_arrowhead = end
    return l

def txt(slide, text, x, y, w, h, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT,
        font='Microsoft YaHei', valign=MSO_ANCHOR.TOP, margin=0.06, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return box

def rich(slide, runs, x, y, w, h, size=20, color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.06):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin); tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    for item in runs:
        r = p.add_run(); r.text = item[0]
        r.font.name = 'Microsoft YaHei'; r.font.size = Pt(item[1] if len(item)>1 else size)
        r.font.bold = item[2] if len(item)>2 else False; r.font.color.rgb = item[3] if len(item)>3 else color
    return box

def title(slide, kicker, heading, page=None):
    txt(slide, kicker.upper(), 0.6, 0.35, 3.7, 0.28, size=10, color=CYAN, bold=True)
    txt(slide, heading, 0.6, 0.68, 11.5, 0.6, size=27, color=NAVY, bold=True)
    line(slide, 0.6, 1.42, 12.7, 1.42, GRID, 1)
    if page is not None: txt(slide, f'{page:02d}', 12.1, 0.42, 0.6, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def footer(slide, text='第一章 运动的描述 · 1.2 时间和位移'):
    line(slide, 0.6, 7.12, 12.7, 7.12, GRID, 1)
    txt(slide, text, 0.6, 7.18, 8, 0.2, size=9, color=MUTED)

def card(slide, x, y, w, h, heading, body='', accent=BLUE, icon=None):
    rect(slide, x, y, w, h, WHITE, GRID, True)
    rect(slide, x, y, 0.08, h, accent, accent, True)
    if icon: txt(slide, icon, x+0.25, y+0.18, 0.5, 0.4, size=22, color=accent, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, heading, x+(0.85 if icon else 0.25), y+0.18, w-(1.05 if icon else 0.5), 0.35, size=16, color=NAVY, bold=True)
    if body: txt(slide, body, x+0.25, y+0.68, w-0.5, h-0.85, size=12.5, color=INK)

def bullet(slide, text, x, y, w, color=INK, size=15, accent=CYAN):
    rect(slide, x, y+0.10, 0.12, 0.12, accent, accent, True)
    txt(slide, text, x+0.25, y, w-0.25, 0.38, size=size, color=color)

def add_image(slide, filename, x, y, w, h=None):
    p = os.path.join(ASSET, filename)
    if os.path.exists(p): slide.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(w), height=Inches(h) if h else None)

def new_slide(bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, bg); return s

# 1 Cover
s = new_slide(NAVY)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 13.333, 0.14, CYAN)
txt(s, '高中物理 · 必修第一册', 0.75, 0.65, 5.4, 0.35, size=15, color=CYAN, bold=True)
txt(s, '1.2 时间和位移', 0.72, 1.45, 8.3, 0.85, size=34, color=WHITE, bold=True)
txt(s, '把“什么时候、在哪里、走了多远”说清楚', 0.78, 2.45, 8.8, 0.45, size=20, color=RGBColor(210, 230, 244))
line(s, 0.8, 3.25, 5.0, 3.25, ORANGE, 4)
txt(s, '开学第一课 · 建模、读图与实验', 0.8, 3.55, 6.2, 0.35, size=15, color=RGBColor(190, 205, 220))
# decorative coordinate motif
line(s, 8.3, 5.9, 12.3, 5.9, CYAN, 2, end='triangle')
line(s, 8.3, 5.9, 8.3, 2.2, CYAN, 2, end='triangle')
for xx, lab in [(8.3,'O'),(9.5,'A'),(11.3,'B')]:
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xx-0.09), Inches(5.81), Inches(0.18), Inches(0.18)).fill.solid();
    sh=s.shapes[-1]; sh.fill.fore_color.rgb=ORANGE; sh.line.color.rgb=ORANGE
    txt(s, lab, xx-0.12, 6.05, 0.25, 0.25, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, '时间轴 × 坐标轴 × 纸带', 8.45, 1.9, 3.5, 0.4, size=16, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
footer(s, '光谷二高 · 高一年级物理')

# 2 Core question
s = new_slide(); title(s, '01 · 情境导入', '如果要准确描述一辆车的运动，我们需要什么？', 2)
add_image(s, '教材批注_introTop.jpg', 0.65, 1.72, 5.25, 4.85)
card(s, 6.25, 1.8, 6.25, 1.25, '先问两个问题', '① 在什么时候？ ② 在哪里？', CYAN, '问')
card(s, 6.25, 3.25, 2.95, 2.1, '时间', '一个瞬间\n还是一段过程？', ORANGE, '时')
card(s, 9.55, 3.25, 2.95, 2.1, '位置', '需要参考系\n和坐标来表达', BLUE, '位')
rect(s, 6.25, 5.8, 6.25, 0.62, PALE, PALE, True)
txt(s, '本节主线：生活语言 → 坐标语言 → 图像语言 → 实验数据', 6.48, 5.98, 5.8, 0.25, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 3 goals roadmap
s = new_slide(); title(s, '02 · 学习目标', '今天下课时，你应该能做到这四件事', 3)
goals = [('辨析', '时刻 vs 时间间隔\n路程 vs 位移', ORANGE, '01'), ('建模', '建立一维坐标系\n读写位置坐标', BLUE, '02'), ('读图', '从 x-t 图像读位置、相遇\n理解斜率与速度', CYAN, '03'), ('测量', '会算打点时间\n会读纸带位移', GREEN, '04')]
for i,(h,b,c,n) in enumerate(goals):
    x=0.75+i*3.1; rect(s,x,1.9,2.75,3.7,WHITE,GRID,True)
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.95), Inches(2.18), Inches(0.85), Inches(0.85)).fill.solid(); o=s.shapes[-1]; o.fill.fore_color.rgb=c; o.line.color.rgb=c
    txt(s,n,x+0.95,2.38,0.85,0.28,size=16,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
    txt(s,h,x+0.25,3.32,2.25,0.38,size=19,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
    txt(s,b,x+0.3,3.9,2.15,0.9,size=14,color=INK,align=PP_ALIGN.CENTER)
    line(s,x+0.45,5.05,x+2.3,5.05,c,3)
txt(s, '一条原则：先规定参考系与正方向，再描述物理量。', 1.1, 6.12, 11.1, 0.4, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 4 moment vs interval
s = new_slide(); title(s, '03 · 时间', '时刻是“点”，时间间隔是“线段”', 4)
txt(s, '时间轴：箭头表示时间增大的方向', 0.85, 1.72, 5, 0.3, size=14, color=MUTED)
line(s, 1.05, 3.1, 11.9, 3.1, NAVY, 2, end='triangle')
for x, lab in [(1.3,'8:00'),(6.05,'8:45'),(10.8,'9:00')]:
    line(s,x,2.86,x,3.34,NAVY,2); txt(s,lab,x-0.4,2.48,0.8,0.25,size=13,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-0.08), Inches(3.02), Inches(0.16), Inches(0.16)).fill.solid(); o=s.shapes[-1]; o.fill.fore_color.rgb=ORANGE; o.line.color.rgb=ORANGE
line(s,1.3,3.8,6.05,3.8,ORANGE,8)
txt(s,'8:00—8:45：45 min（时间间隔）',2.3,4.05,3.2,0.3,size=14,color=ORANGE,bold=True,align=PP_ALIGN.CENTER)
card(s, 0.85, 5.0, 5.35, 1.3, '时刻', '某一瞬间；状态量；没有“长短”\n例：8时上课、8时45分下课', BLUE, '点')
card(s, 6.8, 5.0, 5.35, 1.3, '时间间隔', '两个时刻之间的一段过程；有长短\n例：从8时到8时45分，共45 min', ORANGE, '线')
footer(s)

# 5 fast check
s = new_slide(); title(s, '04 · 课堂辨析', '把日常表达“翻译”成物理语言', 5)
items=[('第4 s末','时刻','与第5 s初是同一时刻',BLUE),('第5 s内','时间间隔','从第4 s初到第5 s末，共1 s',ORANGE),('前3 s内','时间间隔','从0时刻到第3 s末',ORANGE),('第5 s末','时刻','时间轴上用一个点表示',BLUE)]
for i,(a,b,c,col) in enumerate(items):
    y=1.75+i*1.12; rect(s,0.85,y,11.65,0.86,PALE,PALE,True)
    txt(s,a,1.15,y+0.2,2.1,0.3,size=18,color=NAVY,bold=True)
    rect(s,3.55,y+0.18,1.65,0.42,col,col,True); txt(s,b,3.55,y+0.26,1.65,0.2,size=12,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
    txt(s,c,5.65,y+0.2,6.25,0.3,size=15,color=INK)
txt(s,'口诀：看到“某时刻 / 初 / 末”先想点；看到“经过 / 用时 / 内”先想线段。', 1.0, 6.35, 11.4, 0.35, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 6 coordinate system
s = new_slide(); title(s, '05 · 空间描述', '位置不是“感觉”，而是坐标系中的一个数', 6)
rect(s,0.8,1.8,6.0,4.65,PALE,PALE,True)
txt(s,'一维坐标系三要素',1.1,2.05,3.2,0.35,size=20,color=NAVY,bold=True)
line(s,1.35,4.4,6.25,4.4,NAVY,2,end='triangle')
for x,lab in [(1.55,'−20 m'),(3.4,'O / 0'),(5.55,'+30 m')]:
    line(s,x,4.17,x,4.63,NAVY,2); txt(s,lab,x-0.45,4.72,0.9,0.25,size=12,color=NAVY,align=PP_ALIGN.CENTER)
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-0.09), Inches(4.31), Inches(0.18), Inches(0.18)).fill.solid(); o=s.shapes[-1]; o.fill.fore_color.rgb=ORANGE; o.line.color.rgb=ORANGE
txt(s,'原点 O',1.0,5.25,1.3,0.3,size=15,color=ORANGE,bold=True,align=PP_ALIGN.CENTER)
txt(s,'正方向 →',2.8,5.25,1.6,0.3,size=15,color=BLUE,bold=True,align=PP_ALIGN.CENTER)
txt(s,'单位长度',4.8,5.25,1.4,0.3,size=15,color=GREEN,bold=True,align=PP_ALIGN.CENTER)
card(s,7.25,1.8,5.1,1.4,'位置 x','物体在某时刻所处的空间点\n用该点的坐标表示',BLUE,'x')
card(s,7.25,3.48,5.1,1.4,'例：岗亭为原点，向东为正','x=+30 m：岗亭以东30 m\nx=−20 m：岗亭以西20 m',ORANGE,'例')
rect(s,7.25,5.25,5.1,0.95,RGBColor(255,248,231),RGBColor(247,221,163),True)
txt(s,'正方向是人为规定的，关键是“先规定，再使用”。',7.55,5.55,4.5,0.3,size=15,color=INK,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 7 route vs displacement
s = new_slide(); title(s, '06 · 路程与位移', '路线可以不同，初末位置决定的位移只有一个', 7)
txt(s,'北京',0.95,2.3,0.8,0.3,size=17,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
txt(s,'重庆',11.4,4.85,0.8,0.3,size=17,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
line(s,1.75,2.5,11.5,5.0,BLUE,3,end='triangle')
line(s,1.75,2.5,4.0,1.75,ORANGE,3)
line(s,4.0,1.75,11.5,5.0,ORANGE,3,end='triangle')
line(s,1.75,2.5,5.0,5.7, GREEN,3)
line(s,5.0,5.7,11.5,5.0, GREEN,3,end='triangle')
txt(s,'直线距离 ≈ 1300 km',4.7,3.1,2.7,0.3,size=14,color=BLUE,bold=True,align=PP_ALIGN.CENTER)
rect(s,1.0,6.0,5.1,0.55,RGBColor(236,247,255),RGBColor(190,220,240),True); txt(s,'位移 Δx：由初位置指向末位置，与路径无关',1.2,6.17,4.7,0.22,size=13,color=BLUE,bold=True,align=PP_ALIGN.CENTER)
rect(s,6.55,6.0,5.1,0.55,RGBColor(255,246,236),RGBColor(246,217,184),True); txt(s,'路程 s：运动轨迹的长度，路径越绕通常越长',6.75,6.17,4.7,0.22,size=13,color=ORANGE,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 8 scalar/vector
s = new_slide(); title(s, '07 · 标量与矢量', '数字之外，还要问一句：方向呢？', 8)
card(s,0.85,1.85,5.45,3.85,'标量 Scalar','只有大小，没有方向\n\n路程 s · 温度 · 质量 · 时间\n\n路程恒有 s≥0',ORANGE,'量')
card(s,7.0,1.85,5.45,3.85,'矢量 Vector','既有大小，又有方向\n\n位移 Δx · 速度 · 力\n\n一维中：正负号表示方向，\n绝对值表示大小',BLUE,'向')
rect(s,1.3,6.15,10.7,0.56,PALE,PALE,True)
txt(s,'Δx = −3 m 不是“负的长度”，而是：大小 3 m，方向沿 x 轴负方向。',1.55,6.32,10.2,0.22,size=15,color=RED,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 9 formula
s = new_slide(); title(s, '08 · 直线位移', '位移 = 末位置坐标 − 初位置坐标', 9)
rect(s,0.85,1.85,7.25,4.7,PALE,PALE,True)
txt(s,'公式',1.2,2.15,1.1,0.3,size=18,color=CYAN,bold=True)
txt(s,'Δx = x₂ − x₁',1.25,2.65,6.1,0.65,size=32,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
line(s,1.5,4.2,7.35,4.2,NAVY,2,end='triangle')
for x,lab,col in [(2.2,'A\nx₁=+5 m',ORANGE),(5.8,'B\nx₂=+2 m',BLUE)]:
    line(s,x,3.9,x,4.5,NAVY,2); o=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-0.11), Inches(4.09), Inches(0.22), Inches(0.22)); o.fill.solid(); o.fill.fore_color.rgb=col; o.line.color.rgb=col
    txt(s,lab,x-0.6,4.72,1.2,0.55,size=14,color=col,bold=True,align=PP_ALIGN.CENTER)
txt(s,'Δx = 2 − 5 = −3 m',1.35,5.75,6.1,0.35,size=20,color=RED,bold=True,align=PP_ALIGN.CENTER)
card(s,8.55,1.85,3.9,1.35,'方向','Δx<0：沿 x 轴负方向',RED,'←')
card(s,8.55,3.5,3.9,1.35,'大小','|Δx|=3 m',GREEN,'| |')
card(s,8.55,5.15,3.9,1.35,'提醒','坐标要带正负号\n顺序不能颠倒',ORANGE,'记')
footer(s)

# 10 practice
s = new_slide(); title(s, '09 · 练一练', '先画坐标轴，再写“末 − 初”', 10)
rect(s,0.85,1.8,11.65,1.1,RGBColor(255,248,231),RGBColor(246,217,184),True)
txt(s,'物体从 x₁ = −4 m 运动到 x₂ = +6 m，求位移及方向。',1.15,2.15,11.05,0.35,size=20,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
line(s,1.25,4.25,11.9,4.25,NAVY,2,end='triangle')
for x,lab,col in [(3.3,'−4 m',ORANGE),(8.95,'+6 m',BLUE)]:
    line(s,x,3.95,x,4.55,NAVY,2); o=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-0.11), Inches(4.14), Inches(0.22), Inches(0.22)); o.fill.solid(); o.fill.fore_color.rgb=col; o.line.color.rgb=col
    txt(s,lab,x-0.5,4.7,1.0,0.3,size=15,color=col,bold=True,align=PP_ALIGN.CENTER)
line(s,3.3,3.42,8.95,3.42,RED,4,end='triangle')
txt(s,'Δx = (+6) − (−4) = +10 m',2.7,5.45,7.8,0.4,size=23,color=RED,bold=True,align=PP_ALIGN.CENTER)
txt(s,'方向：沿 x 轴正方向；大小：10 m',3.2,6.1,6.9,0.3,size=16,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 11 x-t graph
s = new_slide(); title(s, '10 · x-t 图像', '把“位置随时间变化”画出来', 11)
rect(s,0.85,1.75,7.25,4.95,PALE,PALE,True)
line(s,1.55,6.0,7.55,6.0,NAVY,2,end='triangle'); line(s,1.55,6.0,1.55,2.15,NAVY,2,end='triangle')
for i in range(1,6): line(s,1.55+i,5.93,1.55+i,6.07,GRID,1)
for i in range(1,4): line(s,1.49,6.0-i,1.61,6.0-i,GRID,1)
txt(s,'t / s',7.35,6.12,0.55,0.22,size=12,color=NAVY,bold=True)
txt(s,'x / m',1.05,2.0,0.55,0.22,size=12,color=NAVY,bold=True)
line(s,1.55,5.35,3.2,4.5,BLUE,3); line(s,3.2,4.5,5.4,3.45,BLUE,3); line(s,5.4,3.45,7.15,4.15,ORANGE,3)
for x,y,lab in [(3.2,4.5,'A'),(5.4,3.45,'B'),(7.15,4.15,'C')]:
    o=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-0.09), Inches(y-0.09), Inches(0.18), Inches(0.18)); o.fill.solid(); o.fill.fore_color.rgb=ORANGE; o.line.color.rgb=ORANGE
    txt(s,lab,x+0.12,y-0.18,0.3,0.25,size=13,color=ORANGE,bold=True)
card(s,8.55,1.8,3.9,1.2,'点 A','某时刻、某位置',BLUE,'A')
card(s,8.55,3.25,3.9,1.2,'斜率 k','k=Δx/Δt，表示速度',CYAN,'k')
card(s,8.55,4.7,3.9,1.2,'拐点 C','运动方向可能发生改变',ORANGE,'↪')
txt(s,'注意：x-t 图像不是物体在空间中的运动轨迹。',8.6,6.18,3.8,0.3,size=13,color=RED,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 12 graph reading
s = new_slide(); title(s, '11 · 读图任务', '四个信息点，读懂一张 x-t 图像', 12)
tasks=[('① 纵截距','t=0 时的初始位置 x₀',BLUE),('② 横截距','到达 x=0 的时刻',ORANGE),('③ 交点','同一时刻、同一位置——相遇',GREEN),('④ 斜率','正负看方向，绝对值看快慢',CYAN)]
for i,(h,b,c) in enumerate(tasks):
    x=0.9+(i%2)*6.05; y=1.95+(i//2)*2.05
    rect(s,x,y,5.45,1.55,WHITE,GRID,True); rect(s,x,y,0.08,1.55,c,c,True)
    txt(s,h,x+0.3,y+0.25,4.7,0.3,size=17,color=NAVY,bold=True)
    txt(s,b,x+0.3,y+0.72,4.8,0.45,size=14,color=INK)
rect(s,1.45,6.2,10.45,0.48,RGBColor(255,247,247),RGBColor(245,200,200),True)
txt(s,'同一时刻比较“位置”，同一物体比较“斜率”。',1.75,6.34,9.85,0.22,size=15,color=RED,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 13 measurement overview
s = new_slide(); title(s, '12 · 测量运动', '用“固定时间间隔 + 位置记录”捕捉运动', 13)
card(s,0.85,1.85,3.55,3.9,'频闪照相','每秒闪光10次\n同一张照片留下多个位置\n\n时间间隔由频率决定',CYAN,'频')
card(s,4.9,1.85,3.55,3.9,'电磁打点计时器','约8 V 交流\n振针 + 复写纸\n\n50 Hz → T=0.02 s',BLUE,'磁')
card(s,8.95,1.85,3.55,3.9,'电火花计时器','220 V 交流\n电火花 + 墨粉纸\n\n50 Hz → T=0.02 s',ORANGE,'火')
txt(s,'共同原理：相同时间间隔，在纸带上连续留下点迹，点间距离反映位移。',1.1,6.25,11.0,0.35,size=17,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 14 timer comparison
s = new_slide(); title(s, '13 · 仪器比较', '两种打点计时器：周期相同，留下点迹的方式不同', 14)
headers=['仪器','电源与频率','打点周期','介质 / 点迹']
xs=[0.85,3.15,6.15,8.3]; ws=[2.0,2.7,1.8,4.2]
for x,w,h in zip(xs,ws,headers): rect(s,x,1.9,w,0.62,NAVY,NAVY,True); txt(s,h,x+0.05,2.09,w-0.1,0.2,size=13,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
rows=[('电磁打点计时器','约8 V 交流，50 Hz','0.02 s','振针在复写纸上连续打点',BLUE),('电火花计时器','220 V 交流，50 Hz','0.02 s','电火花在墨粉纸上留下点迹',ORANGE)]
for r,(a,b,c,d,col) in enumerate(rows):
    y=2.65+r*1.2
    for x,w,val in zip(xs,ws,[a,b,c,d]): rect(s,x,y,w,0.95,WHITE,GRID,False); txt(s,val,x+0.08,y+0.28,w-0.16,0.35,size=13,color=INK,bold=(x==0.85),align=PP_ALIGN.CENTER)
    rect(s,0.85,y,0.08,0.95,col,col)
rect(s,0.85,5.45,11.65,0.85,RGBColor(236,247,255),RGBColor(190,220,240),True)
txt(s,'T = 1/f = 1/50 s = 0.02 s',1.25,5.72,10.8,0.3,size=23,color=BLUE,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 15 experiment flow
s = new_slide(); title(s, '14 · 纸带实验', '规范操作：先通电，再拉动', 15)
steps=[('1','固定仪器','穿好纸带',BLUE),('2','安装纸带','压在复写纸/墨粉纸下',CYAN),('3','接通电源','稳定后水平拉动',ORANGE),('4','断电取带','选清晰点作起点',GREEN),('5','测量记录','尺量位移，表格记数据',RED)]
for i,(n,h,b,c) in enumerate(steps):
    x=0.75+i*2.45
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.7), Inches(2.05), Inches(0.9), Inches(0.9)).fill.solid(); o=s.shapes[-1]; o.fill.fore_color.rgb=c; o.line.color.rgb=c
    txt(s,n,x+0.7,2.29,0.9,0.3,size=22,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
    if i<4: line(s,x+1.6,2.5,x+2.25,2.5,GRID,3,end='triangle')
    txt(s,h,x+0.1,3.2,2.2,0.35,size=16,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
    txt(s,b,x+0.1,3.72,2.2,0.65,size=13,color=INK,align=PP_ALIGN.CENTER)
rect(s,1.15,5.25,11.0,0.9,RGBColor(255,247,247),RGBColor(245,200,200),True)
txt(s,'关键提醒：电源应先接通再拉纸带；纸带要水平、点迹要清晰。',1.45,5.54,10.4,0.3,size=18,color=RED,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 16 counting points
s = new_slide(); title(s, '15 · 数据处理', '数出 n 个点，只有 n−1 个时间间隔', 16)
rect(s,0.85,1.85,7.25,4.65,PALE,PALE,True)
txt(s,'纸带点迹（示意）',1.15,2.15,2.5,0.3,size=16,color=NAVY,bold=True)
for i in range(8):
    x=1.45+i*0.78; y=4.15-0.05*i
    o=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.16), Inches(0.16)); o.fill.solid(); o.fill.fore_color.rgb=ORANGE; o.line.color.rgb=ORANGE
    txt(s,str(i+1),x-0.05,4.45,0.25,0.2,size=11,color=MUTED,align=PP_ALIGN.CENTER)
for i in range(7): line(s,1.61+i*0.78,4.22,2.23+i*0.78,4.22,BLUE,1)
txt(s,'8 个点 → 7 个间隔',2.2,5.0,3.9,0.35,size=19,color=BLUE,bold=True,align=PP_ALIGN.CENTER)
card(s,8.55,1.85,3.9,1.3,'时间','t=(n−1)T',BLUE,'t')
card(s,8.55,3.4,3.9,1.3,'本例','(8−1)×0.02=0.14 s',ORANGE,'算')
card(s,8.55,4.95,3.9,1.3,'位移','用刻度尺量起点到终点',GREEN,'x')
footer(s)

# 17 data table & errors
s = new_slide(); title(s, '16 · 实验记录', '原始数据要如实记录，异常现象要能解释', 17)
headers=['终点点号 n','间隔数 n−1','时间 t/s','位移 x₀/m','备注']
xs=[1.0,3.0,5.05,7.1,9.15]; ws=[1.75,1.85,1.75,1.75,2.9]
for x,w,h in zip(xs,ws,headers): rect(s,x,1.85,w,0.55,NAVY,NAVY,True); txt(s,h,x+0.03,2.03,w-0.06,0.18,size=12,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
for r in range(3):
    y=2.43+r*0.62
    for x,w in zip(xs,ws): rect(s,x,y,w,0.62,WHITE,GRID,False)
    txt(s,str(r+2),1.0,y+0.18,1.75,0.2,size=13,color=INK,align=PP_ALIGN.CENTER)
    txt(s,str(r+1),3.0,y+0.18,1.85,0.2,size=13,color=INK,align=PP_ALIGN.CENTER)
    txt(s,f'{(r+1)*0.02:.2f}',5.05,y+0.18,1.75,0.2,size=13,color=INK,align=PP_ALIGN.CENTER)
txt(s,'点迹模糊？',1.0,5.05,1.7,0.35,size=18,color=RED,bold=True)
for i,t in enumerate(['电压过低','振针过高','纸带过松','纸带未压好']): bullet(s,t,2.8+(i%2)*4.5,4.98+(i//2)*0.62,3.6,size=14,accent=RED)
footer(s)

# 18 summary map
s = new_slide(); title(s, '17 · 知识结构', '从“时间”到“位移”，再到“图像与实验”', 18)
nodes=[('时间','时刻 · 时间间隔',0.9,2.2,ORANGE),('空间','坐标系 · 位置 x',3.45,2.2,BLUE),('变化','路程 s · 位移 Δx',6.0,2.2,CYAN),('图像','x-t · 斜率=速度',8.55,2.2,GREEN),('测量','频闪 · 打点纸带',4.7,4.85,RED)]
for h,b,x,y,c in nodes:
    rect(s,x,y,2.1,1.25,WHITE,c,True); txt(s,h,x+0.15,y+0.18,1.8,0.28,size=17,color=c,bold=True,align=PP_ALIGN.CENTER); txt(s,b,x+0.12,y+0.62,1.86,0.3,size=12,color=INK,align=PP_ALIGN.CENTER)
line(s,2.98,2.82,3.42,2.82,GRID,2,end='triangle'); line(s,5.53,2.82,5.97,2.82,GRID,2,end='triangle'); line(s,8.08,2.82,8.52,2.82,GRID,2,end='triangle')
line(s,7.0,3.47,5.75,4.8,GRID,2,end='triangle'); line(s,4.95,3.47,5.3,4.8,GRID,2,end='triangle')
rect(s,1.1,6.32,11.0,0.48,PALE,PALE,True); txt(s,'一条原则：先规定参考系和正方向，再描述物理量。',1.35,6.45,10.5,0.2,size=15,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 19 exit ticket
s = new_slide(); title(s, '18 · 出口条', '离开教室前，请完成这 5 个小问题', 19)
qs=['① “第6 s末”和“第7 s初”是同一时刻吗？','② x₁=−4 m 到 x₂=+6 m，位移是多少？','③ 绕一周回到原点，路程与位移分别怎样？','④ x-t 图像交点、斜率分别表示什么？','⑤ 数到第8个点，50 Hz，经历多长时间？']
for i,q in enumerate(qs):
    y=1.78+i*0.85; rect(s,0.95,y,11.45,0.58,PALE,PALE,True); txt(s,q,1.25,y+0.16,10.85,0.25,size=15,color=INK)
rect(s,1.35,6.2,10.65,0.5,RGBColor(236,247,255),RGBColor(190,220,240),True); txt(s,'答案关键词：是；+10 m；s>0、Δx=0；相遇、速度；0.14 s',1.55,6.34,10.25,0.2,size=13,color=BLUE,bold=True,align=PP_ALIGN.CENTER)
footer(s)

# 20 extension
s = new_slide(NAVY)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,0,13.333,0.14,ORANGE)
txt(s,'课后延伸 · 把生活拍成一张 x-t 图',0.75,0.7,8.8,0.4,size=20,color=YELLOW,bold=True)
txt(s,'用手机固定时间间隔拍摄小球或玩具车的运动',0.75,1.45,9.4,0.6,size=28,color=WHITE,bold=True)
for i in range(6):
    x=1.0+i*1.8; y=4.2-0.25*math.sin(i)
    o=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.28), Inches(0.28)); o.fill.solid(); o.fill.fore_color.rgb=CYAN; o.line.color.rgb=CYAN
    if i<5: line(s,x+0.28,y+0.14,x+1.78,y+0.14,RGBColor(93,151,186),2,end='triangle')
line(s,1.0,5.6,11.9,5.6,RGBColor(93,151,186),2,end='triangle'); line(s,1.0,5.6,1.0,3.15,RGBColor(93,151,186),2,end='triangle')
txt(s,'下节课：用纸带实验数据，与照片数据“对照验证”。',1.0,6.35,11.2,0.35,size=17,color=RGBColor(205,225,240),align=PP_ALIGN.CENTER)
footer(s, '光谷二高 · 高一年级物理')

prs.save(OUT)
print(OUT)
